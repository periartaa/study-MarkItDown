import os
import pdfplumber
import docx
import openpyxl
from pptx import Presentation

def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    if ext == ".pdf":
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                content += page.extract_text() + "\n"

    elif ext == ".docx":
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            content += para.text + "\n"

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"

    elif ext == ".xlsx":
        wb = openpyxl.load_workbook(file_path)
        for sheet in wb.worksheets:
            content += f"# {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                content += ' | '.join([str(cell) if cell is not None else "" for cell in row]) + "\n"

    else:
        content = "Format file tidak didukung."

    return content

if __name__ == "__main__":
    file_path = input("Masukkan path file yang ingin dibaca: ").strip()
    
    # Periksa apakah file ada
    if not os.path.exists(file_path):
        print("File tidak ditemukan!")
    else:
        try:
            content = read_file(file_path)
            print("\nIsi file:\n")
            print(content)
        except Exception as e:
            print(f"Terjadi error saat membaca file: {e}")